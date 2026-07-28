#!/usr/bin/env python3
"""
Belge Yükleme Scripti - PDF, PPTX, DOCX dosyalarını Pinecone'a yükler.

Kullanım:
    python scripts/ingest.py --path data/notes --category ders_notu
    python scripts/ingest.py --path data/psychology --category psikoloji
    python scripts/ingest.py --path data/clinical --category klinik
"""

import os
import sys
import argparse
from pathlib import Path
import hashlib

sys.path.insert(0, str(Path(__file__).parent.parent))

from dotenv import load_dotenv
load_dotenv()

from pinecone import Pinecone, ServerlessSpec
from pypdf import PdfReader
from pptx import Presentation
from docx import Document

from app.config import get_settings
from app.rag.embeddings import get_embedding
from app.prompts import DocumentCategory


def extract_text_from_pdf(file_path: str) -> list[dict]:
    """PDF'den metin çıkar."""
    pages = []
    try:
        reader = PdfReader(file_path)
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append({
                    "text": text.strip(),
                    "page": i + 1,
                    "source": Path(file_path).name,
                })
    except Exception as e:
        print(f"   ❌ PDF okuma hatası: {e}")
    return pages


def extract_text_from_pptx(file_path: str) -> list[dict]:
    """PPTX'den metin çıkar."""
    slides = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
            if texts:
                slides.append({
                    "text": "\n".join(texts),
                    "page": i + 1,
                    "source": Path(file_path).name,
                })
    except Exception as e:
        print(f"   ❌ PPTX okuma hatası: {e}")
    return slides


def extract_text_from_docx(file_path: str) -> list[dict]:
    """DOCX'den metin çıkar."""
    paragraphs = []
    try:
        doc = Document(file_path)
        full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        if full_text:
            paragraphs.append({
                "text": full_text,
                "page": 1,
                "source": Path(file_path).name,
            })
    except Exception as e:
        print(f"   ❌ DOCX okuma hatası: {e}")
    return paragraphs


def load_documents(directory: str) -> list[dict]:
    """Dizindeki tüm desteklenen dosyaları yükler."""
    documents = []
    path = Path(directory)
    
    if not path.exists():
        print(f"❌ Dizin bulunamadı: {directory}")
        return documents
    
    handlers = {
        ".pdf": extract_text_from_pdf,
        ".pptx": extract_text_from_pptx,
        ".ppt": extract_text_from_pptx,
        ".docx": extract_text_from_docx,
        ".doc": extract_text_from_docx,
    }
    
    for ext, handler in handlers.items():
        for file_path in path.glob(f"**/*{ext}"):
            print(f"📄 Yükleniyor: {file_path.name}")
            docs = handler(str(file_path))
            documents.extend(docs)
            print(f"   ✓ {len(docs)} parça yüklendi")
    
    return documents


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> list[str]:
    """Metni parçalara böler."""
    if len(text) <= chunk_size:
        return [text]
    
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        
        if end < len(text):
            last_space = chunk.rfind(' ')
            if last_space > chunk_size // 2:
                end = start + last_space
                chunk = text[start:end]
        
        chunks.append(chunk.strip())
        start = end - overlap
    
    return chunks


def ensure_index_exists(settings):
    """Pinecone index'inin var olduğundan emin ol."""
    pc = Pinecone(api_key=settings.pinecone_api_key)
    
    existing_indexes = [idx.name for idx in pc.list_indexes()]
    
    if settings.pinecone_index_name not in existing_indexes:
        print(f"📌 Yeni index oluşturuluyor: {settings.pinecone_index_name}")
        pc.create_index(
            name=settings.pinecone_index_name,
            dimension=768,
            metric="cosine",
            spec=ServerlessSpec(
                cloud="aws",
                region="us-east-1"
            )
        )
        print("   ✓ Index oluşturuldu")
    else:
        print(f"📌 Index mevcut: {settings.pinecone_index_name}")
    
    return pc.Index(settings.pinecone_index_name)


def ingest_documents(directory: str, category: str):
    """Ana yükleme fonksiyonu."""
    settings = get_settings()
    
    print(f"\n{'='*50}")
    print(f"🏥 Canan Hemşirelik Asistanı - Belge Yükleme")
    print(f"{'='*50}")
    print(f"📁 Dizin: {directory}")
    print(f"🏷️  Kategori: {category}")
    print(f"{'='*50}\n")
    
    print("1️⃣ Belgeler yükleniyor...")
    documents = load_documents(directory)
    
    if not documents:
        print("❌ Yüklenecek belge bulunamadı!")
        return
    
    print("\n2️⃣ Belgeler parçalanıyor...")
    chunks = []
    for doc in documents:
        text_chunks = chunk_text(
            doc["text"],
            chunk_size=settings.chunk_size,
            overlap=settings.chunk_overlap,
        )
        for chunk in text_chunks:
            chunks.append({
                "text": chunk,
                "source": doc["source"],
                "page": doc["page"],
                "category": category,
            })
    
    print(f"📦 {len(documents)} belge → {len(chunks)} parça")
    
    print("\n3️⃣ Pinecone bağlantısı kontrol ediliyor...")
    index = ensure_index_exists(settings)
    
    print("\n4️⃣ Embedding'ler oluşturuluyor ve yükleniyor...")
    
    vectors = []
    for i, chunk in enumerate(chunks):
        print(f"   İşleniyor: {i+1}/{len(chunks)}", end="\r")
        
        chunk_id = hashlib.md5(
            f"{chunk['source']}_{chunk['page']}_{i}".encode()
        ).hexdigest()
        
        embedding = get_embedding(chunk["text"])
        
        vectors.append({
            "id": chunk_id,
            "values": embedding,
            "metadata": {
                "text": chunk["text"],
                "source": chunk["source"],
                "page": str(chunk["page"]),
                "category": chunk["category"],
            }
        })
        
        if len(vectors) >= 50:
            index.upsert(vectors=vectors)
            vectors = []
    
    if vectors:
        index.upsert(vectors=vectors)
    
    print(f"\n\n{'='*50}")
    print(f"✅ TAMAMLANDI!")
    print(f"📊 Toplam {len(chunks)} parça Pinecone'a yüklendi")
    print(f"🏷️  Kategori: {category}")
    print(f"{'='*50}\n")


def main():
    parser = argparse.ArgumentParser(
        description="Hemşirelik belgelerini Pinecone'a yükle"
    )
    parser.add_argument(
        "--path",
        type=str,
        required=True,
        help="Belgelerin bulunduğu dizin (örn: data/notes)"
    )
    parser.add_argument(
        "--category",
        type=str,
        required=True,
        choices=[c.value for c in DocumentCategory],
        help="Belge kategorisi: ders_notu, psikoloji, klinik, farmakoloji, genel"
    )
    
    args = parser.parse_args()
    ingest_documents(args.path, args.category)


if __name__ == "__main__":
    main()
